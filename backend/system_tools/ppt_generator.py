# backend/system_tools/ppt_generator.py
"""
PPT 生成工具：根据结构化幻灯片数据生成 .pptx 文件。

支持：
- 4 种页面布局：cover（封面）、toc（目录）、content（内容）、ending（结尾）
- 3 套配色主题：blue（简约蓝）、warm（活力橙）、clean（极简白）
- 基础 Markdown 解析：**加粗**、- 列表、1. 编号
- 图片嵌入：image_url（直链）或 image_query（Unsplash 免费图库自动搜索）
- 生成 .pptx → 返回下载链接 + 推送前端预览标记

Agent 执行流程：
  用户输入 → system_todo 规划大纲 → LLM 生成 slides JSON
  → system_generate_pptx → 前端渲染幻灯片卡片 + 下载链接
"""
import json
import os
import uuid
import hashlib
import tempfile
import urllib.request
import urllib.error
import ssl
from pathlib import Path
from typing import Any, Dict, List, Optional

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

from config_loader import config
import backend

# ── 预览事件队列 ──
_ppt_preview_events: List[str] = []


def push_ppt_preview(marker: str):
    """追加一条 PPT 预览标记到事件队列。"""
    _ppt_preview_events.append(marker)


def pop_ppt_preview_events() -> List[str]:
    """取出并清空当前累积的所有 PPT 预览事件标记。"""
    global _ppt_preview_events
    events = _ppt_preview_events
    _ppt_preview_events = []
    return events


# ── 主题配置 ──
THEMES = {
    "blue": {
        "name": "简约蓝",
        "bg": RGBColor(0xFF, 0xFF, 0xFF),
        "title_color": RGBColor(0x25, 0x63, 0xEB),
        "text_color": RGBColor(0x1E, 0x29, 0x3B),
        "accent": RGBColor(0x25, 0x63, 0xEB),
        "accent_light": RGBColor(0xDB, 0xEA, 0xFE),
        "footer_color": RGBColor(0x94, 0xA3, 0xB8),
    },
    "warm": {
        "name": "活力橙",
        "bg": RGBColor(0xFF, 0xFB, 0xEB),
        "title_color": RGBColor(0xEA, 0x58, 0x0C),
        "text_color": RGBColor(0x44, 0x40, 0x3C),
        "accent": RGBColor(0xEA, 0x58, 0x0C),
        "accent_light": RGBColor(0xFE, 0xD7, 0xAA),
        "footer_color": RGBColor(0xA8, 0xA2, 0x9E),
    },
    "clean": {
        "name": "极简白",
        "bg": RGBColor(0xFF, 0xFF, 0xFF),
        "title_color": RGBColor(0x0F, 0x17, 0x2A),
        "text_color": RGBColor(0x47, 0x55, 0x69),
        "accent": RGBColor(0x47, 0x55, 0x69),
        "accent_light": RGBColor(0xF1, 0xF5, 0xF9),
        "footer_color": RGBColor(0x94, 0xA3, 0xB8),
    },
}

DEFAULT_THEME = "blue"

# ── 字体配置 ──
FONT_FAMILY = "Microsoft YaHei"

# ── 幻灯片尺寸（16:9 宽屏） ──
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# ── 图片下载超时（秒）──
IMAGE_DOWNLOAD_TIMEOUT = 15

# ── 图片缓存目录 ──
_IMAGE_CACHE_DIR = None


def _get_image_cache_dir() -> str:
    """获取图片缓存目录（惰性创建）。"""
    global _IMAGE_CACHE_DIR
    if _IMAGE_CACHE_DIR is None:
        _IMAGE_CACHE_DIR = os.path.join(str(config.generate_dir), ".image_cache")
        os.makedirs(_IMAGE_CACHE_DIR, exist_ok=True)
    return _IMAGE_CACHE_DIR


def _get_font_family() -> str:
    return FONT_FAMILY


def _apply_theme_colors(theme_name: str) -> Dict[str, Any]:
    return THEMES.get(theme_name, THEMES[DEFAULT_THEME])


# ── 图片下载与处理 ──

def _download_image(url: str) -> Optional[str]:
    """从 URL 下载图片到本地缓存，返回本地路径。失败返回 None。"""
    try:
        # 生成缓存文件名（URL 哈希）
        url_hash = hashlib.md5(url.encode()).hexdigest()[:12]
        # 推断扩展名
        ext = ".jpg"
        url_lower = url.lower()
        if ".png" in url_lower:
            ext = ".png"
        elif ".webp" in url_lower:
            ext = ".webp"
        elif ".gif" in url_lower:
            ext = ".gif"

        cache_path = os.path.join(_get_image_cache_dir(), f"{url_hash}{ext}")

        # 命中缓存直接返回
        if os.path.exists(cache_path) and os.path.getsize(cache_path) > 0:
            return cache_path

        # 下载（忽略 SSL 验证以兼容更多图源）
        ctx = ssl.create_default_context()
        ctx.check_hostname = False
        ctx.verify_mode = ssl.CERT_NONE

        req = urllib.request.Request(url, headers={
            "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36"
        })
        with urllib.request.urlopen(req, timeout=IMAGE_DOWNLOAD_TIMEOUT, context=ctx) as resp:
            if resp.status not in (200, 301, 302):
                return None
            data = resp.read()
            if len(data) < 1024:  # 太小，不是有效图片
                return None
            with open(cache_path, "wb") as f:
                f.write(data)
            return cache_path
    except Exception:
        return None


def _resolve_image(slide_data: Dict) -> Optional[str]:
    """解析幻灯片的配图：优先 image_url，其次 image_query。

    对于 image_query，使用多级回退策略搜索免费图片：
    1. Picsum Photos（基于 query 哈希的稳定随机图片）
    2. 若全部失败，返回 None（幻灯片无图但不影响整体生成）

    返回本地图片路径，失败返回 None。
    """
    # 预取缓存命中（由 generate_pptx 并行预下载注入）
    cached = slide_data.get("_cached_image")
    if cached:
        return cached

    # 优先直接 URL
    image_url = slide_data.get("image_url", "").strip()
    if image_url:
        return _download_image(image_url)

    # 其次关键词搜索
    image_query = slide_data.get("image_query", "").strip()
    if image_query:
        from urllib.parse import quote
        query_encoded = quote(image_query, safe='')

        # 策略 1：用 query 哈希作为 Picsum seed，相同关键词总是返回同一张图
        query_hash = abs(hash(image_query)) % 1000
        picsum_url = f"https://picsum.photos/seed/{query_hash}/800/600"
        path = _download_image(picsum_url)
        if path:
            return path

    return None


# ── Markdown 解析 ──

def _parse_markdown_runs(text: str, paragraph, theme: Dict):
    """将基础 Markdown 文本转换为 python-pptx 的富文本 runs。"""
    import re

    font_family = _get_font_family()
    text_color = theme["text_color"]
    text = text.strip()

    bullet_match = re.match(r'^(\s*)[-*]\s+(.+)$', text)
    numbered_match = re.match(r'^(\s*)\d+[.)]\s+(.+)$', text)

    if bullet_match:
        paragraph.level = min(len(bullet_match.group(1)) // 2, 8)
        text = bullet_match.group(2)
    elif numbered_match:
        paragraph.level = min(len(numbered_match.group(1)) // 2, 8)
        text = numbered_match.group(2)

    segments = re.split(r'(\*\*.*?\*\*|\*[^*]+\*)', text)

    for seg in segments:
        if not seg:
            continue
        if seg.startswith('**') and seg.endswith('**'):
            run = paragraph.add_run()
            run.text = seg[2:-2]
            run.font.bold = True
            run.font.size = Pt(18)
            run.font.color.rgb = text_color
            run.font.name = font_family
        elif seg.startswith('*') and seg.endswith('*') and len(seg) > 2:
            run = paragraph.add_run()
            run.text = seg[1:-1]
            run.font.italic = True
            run.font.size = Pt(18)
            run.font.color.rgb = text_color
            run.font.name = font_family
        else:
            run = paragraph.add_run()
            run.text = seg
            run.font.size = Pt(18)
            run.font.color.rgb = text_color
            run.font.name = font_family


# ── 布局函数 ──

def _add_cover_slide(prs, slide_data: Dict, theme: Dict, title: str = ""):
    """添加封面页。"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = theme["bg"]

    # 顶部装饰条
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.08)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = theme["accent"]
    shape.line.fill.background()

    # 主标题
    slide_title = slide_data.get("title", title or "演示文稿")
    txBox = slide.shapes.add_textbox(Inches(1.5), Inches(2.2), Inches(10.3), Inches(1.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = slide_title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = theme["title_color"]
    p.font.name = _get_font_family()
    p.alignment = PP_ALIGN.CENTER

    subtitle = slide_data.get("subtitle", "")
    if subtitle:
        txBox2 = slide.shapes.add_textbox(Inches(1.5), Inches(4.2), Inches(10.3), Inches(1.0))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = subtitle
        p2.font.size = Pt(22)
        p2.font.color.rgb = theme["text_color"]
        p2.font.name = _get_font_family()
        p2.alignment = PP_ALIGN.CENTER

    speaker = slide_data.get("speaker", "")
    if speaker:
        txBox3 = slide.shapes.add_textbox(Inches(1.5), Inches(6.2), Inches(10.3), Inches(0.6))
        tf3 = txBox3.text_frame
        p3 = tf3.paragraphs[0]
        p3.text = speaker
        p3.font.size = Pt(14)
        p3.font.color.rgb = theme["footer_color"]
        p3.font.name = _get_font_family()
        p3.alignment = PP_ALIGN.CENTER


def _add_toc_slide(prs, slide_data: Dict, theme: Dict, index: int):
    """添加目录页。"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = theme["bg"]

    slide_title = slide_data.get("title", "目录")
    txBox = slide.shapes.add_textbox(Inches(1.0), Inches(0.6), Inches(11.3), Inches(0.9))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = slide_title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = theme["title_color"]
    p.font.name = _get_font_family()

    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(1.7), Inches(11.3), Inches(0.03)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = theme["accent"]
    shape.line.fill.background()

    items = slide_data.get("items", [])
    if not items:
        items = slide_data.get("bullets", [])

    y_start = 2.2
    for i, item in enumerate(items, 1):
        txNum = slide.shapes.add_textbox(Inches(1.2), Inches(y_start + i * 0.7), Inches(0.8), Inches(0.5))
        tn = txNum.text_frame
        tnp = tn.paragraphs[0]
        tnp.text = f"{i:02d}"
        tnp.font.size = Pt(24)
        tnp.font.bold = True
        tnp.font.color.rgb = theme["accent"]
        tnp.font.name = _get_font_family()

        txItem = slide.shapes.add_textbox(Inches(2.2), Inches(y_start + i * 0.7), Inches(9.5), Inches(0.5))
        ti = txItem.text_frame
        tip = ti.paragraphs[0]
        tip.text = item
        tip.font.size = Pt(20)
        tip.font.color.rgb = theme["text_color"]
        tip.font.name = _get_font_family()


def _add_content_slide(prs, slide_data: Dict, theme: Dict, index: int):
    """添加内容页，支持右侧配图（图文并茂）。"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = theme["bg"]

    # ── 解析配图 ──
    image_path = _resolve_image(slide_data)
    has_image = image_path and os.path.exists(image_path)

    # ── 根据是否有图片调整正文区域宽度 ──
    if has_image:
        text_width = Inches(6.5)   # 左侧文本区域
        text_left = Inches(0.8)
        image_left = Inches(7.8)
        image_top = Inches(1.8)
        image_width = Inches(4.8)
        image_height = Inches(4.5)
    else:
        text_width = Inches(11.3)
        text_left = Inches(1.0)

    # 页面标题
    slide_title = slide_data.get("title", f"第 {index + 1} 页")
    txBox = slide.shapes.add_textbox(Inches(1.0), Inches(0.5), Inches(11.3), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = slide_title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = theme["title_color"]
    p.font.name = _get_font_family()

    # 标题下划线
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(1.4), Inches(2.5), Inches(0.05)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = theme["accent"]
    shape.line.fill.background()

    # 正文区域
    content = slide_data.get("content", "")
    bullets = slide_data.get("bullets", [])

    txBody = slide.shapes.add_textbox(text_left, Inches(1.8), text_width, Inches(5.0))
    tf_body = txBody.text_frame
    tf_body.word_wrap = True

    lines = []
    if bullets:
        lines = [f"- {b}" for b in bullets]
    elif content:
        lines = content.strip().split('\n')

    if not lines:
        lines = [""]

    for i, line in enumerate(lines):
        if i == 0:
            para = tf_body.paragraphs[0]
        else:
            para = tf_body.add_paragraph()
        _parse_markdown_runs(line, para, theme)

    # 嵌入配图
    if has_image:
        try:
            pic = slide.shapes.add_picture(
                image_path, image_left, image_top, image_width, image_height
            )
            # 图片边框
            pic.line.color.rgb = theme["accent_light"]
            pic.line.width = Pt(1)
        except Exception:
            pass  # 图片嵌入失败不影响整体

    # 页码
    txPage = slide.shapes.add_textbox(Inches(11.5), Inches(7.0), Inches(1.5), Inches(0.4))
    tp = txPage.text_frame
    tpp = tp.paragraphs[0]
    tpp.text = str(index + 1)
    tpp.font.size = Pt(11)
    tpp.font.color.rgb = theme["footer_color"]
    tpp.font.name = _get_font_family()
    tpp.alignment = PP_ALIGN.RIGHT


def _add_ending_slide(prs, slide_data: Dict, theme: Dict, index: int):
    """添加结尾页。"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = theme["bg"]

    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.08)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = theme["accent"]
    shape.line.fill.background()

    slide_title = slide_data.get("title", "感谢聆听")
    txBox = slide.shapes.add_textbox(Inches(1.5), Inches(2.5), Inches(10.3), Inches(1.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = slide_title
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = theme["title_color"]
    p.font.name = _get_font_family()
    p.alignment = PP_ALIGN.CENTER

    subtitle = slide_data.get("subtitle", "")
    if subtitle:
        txBox2 = slide.shapes.add_textbox(Inches(1.5), Inches(4.2), Inches(10.3), Inches(0.8))
        tf2 = txBox2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = subtitle
        p2.font.size = Pt(20)
        p2.font.color.rgb = theme["text_color"]
        p2.font.name = _get_font_family()
        p2.alignment = PP_ALIGN.CENTER


_LAYOUT_HANDLERS = {
    "cover": _add_cover_slide,
    "toc": _add_toc_slide,
    "content": _add_content_slide,
    "ending": _add_ending_slide,
}


async def generate_pptx(
    slides: list = None,
    theme: str = "blue",
    filename: str = "",
    title: str = "",
    **_kwargs,
) -> Dict[str, Any]:
    """
    根据结构化幻灯片数据生成 .pptx 演示文稿文件。

    Args:
        slides: 幻灯片数组，每项 {title, type?, content?, bullets?, subtitle?,
                image_query?, image_url?, speaker?, items?}
        theme: blue / warm / clean（默认 blue）
        filename: 输出文件名（不含扩展名），省略则自动生成
        title: 演示文稿总标题

    Returns:
        {success, path, download_url, slide_count, slides_preview, image_count}
    """
    slides = slides or []
    if not isinstance(slides, list) or not slides:
        return {"success": False, "error": "slides 必须是非空数组"}

    theme_name = theme if theme in THEMES else DEFAULT_THEME
    theme_colors = _apply_theme_colors(theme_name)

    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # ── 并行预下载所有图片（线程池） ──
    # 原逻辑每张图片调用 _resolve_image 两次（handler + 统计），
    # 且 urllib 同步串行下载。改为一次性并行预取，大幅提速。
    from concurrent.futures import ThreadPoolExecutor as _ThreadPoolExecutor
    _image_specs = []  # [(slide_index, slide_data)]
    for i, sd in enumerate(slides):
        if isinstance(sd, dict) and sd.get("type", "content") in (None, "content"):
            if sd.get("image_url") or sd.get("image_query"):
                _image_specs.append((i, sd))

    _image_results = {}  # slide_index → local_path or None
    if _image_specs:
        with _ThreadPoolExecutor(max_workers=6) as _pool:
            _futures = {_pool.submit(_resolve_image, sd): i for i, sd in _image_specs}
            for _fut in _futures:
                try:
                    _image_results[_futures[_fut]] = _fut.result(timeout=IMAGE_DOWNLOAD_TIMEOUT)
                except Exception:
                    _image_results[_futures[_fut]] = None

    # 注入图片路径到 slide_data，避免 handler 内部二次调用 _resolve_image
    for i, sd in enumerate(slides):
        if isinstance(sd, dict) and i in _image_results and _image_results[i]:
            sd = sd.copy()
            sd["_cached_image"] = _image_results[i]
            slides[i] = sd

    preview_slides = []
    image_count = 0

    for i, slide_data in enumerate(slides):
        if not isinstance(slide_data, dict):
            continue
        slide_type = slide_data.get("type", "content")
        handler = _LAYOUT_HANDLERS.get(slide_type, _add_content_slide)
        handler(prs, slide_data, theme_colors, i)

        # 统计图片（复用预取结果，无需再次下载）
        cached = slide_data.get("_cached_image")
        if cached:
            image_count += 1

        # 预览数据
        raw_content = slide_data.get("content", "")
        bullets = slide_data.get("bullets", [])
        if bullets:
            content_preview = " · ".join(bullets[:3])
        elif raw_content:
            content_preview = raw_content[:120].replace('\n', ' ')
        else:
            content_preview = ""

        preview_slides.append({
            "title": slide_data.get("title", ""),
            "content_preview": content_preview,
            "type": slide_type,
            "has_image": bool(slide_data.get("image_url") or slide_data.get("image_query")),
        })

    # ── 输出路径 ──
    output_dir = config.generate_dir
    os.makedirs(output_dir, exist_ok=True)

    safe_filename = filename.strip() if filename and filename.strip() else ""
    if not safe_filename:
        safe_filename = f"presentation_{uuid.uuid4().hex[:8]}"
    safe_filename = "".join(c for c in safe_filename if c.isalnum() or c in "._-（）")
    if not safe_filename:
        safe_filename = f"presentation_{uuid.uuid4().hex[:8]}"

    pptx_path = os.path.join(str(output_dir), f"{safe_filename}.pptx")

    counter = 1
    while os.path.exists(pptx_path):
        pptx_path = os.path.join(str(output_dir), f"{safe_filename}_{counter}.pptx")
        counter += 1

    # ── 保存 ──
    try:
        prs.save(pptx_path)
    except Exception as e:
        return {"success": False, "error": f"保存 PPTX 文件失败：{e}"}

    # ── 下载链接 ──
    rel_path = os.path.relpath(pptx_path, str(config.generate_dir))
    download_url = f"/files/generate/{rel_path.replace(os.sep, '/')}"

    # ── 预览数据 ──
    preview_data = {
        "file": download_url,
        "filename": os.path.basename(pptx_path),
        "title": title or safe_filename,
        "slide_count": len(preview_slides),
        "theme": theme_name,
        "theme_label": theme_colors["name"],
        "image_count": image_count,
        "slides": preview_slides,
    }

    # ── 推送预览标记 ──
    import base64
    marker_json = json.dumps(preview_data, ensure_ascii=False)
    encoded = base64.b64encode(marker_json.encode('utf-8')).decode('ascii')
    marker = f"<!--ppt_preview:{encoded}-->"
    push_ppt_preview(marker)

    img_info = f"，含 {image_count} 张配图" if image_count > 0 else ""
    return {
        "success": True,
        "path": pptx_path,
        "download_url": download_url,
        "filename": os.path.basename(pptx_path),
        "slide_count": len(preview_slides),
        "image_count": image_count,
        "slides_preview": preview_slides,
        "message": (
            f"已生成 {len(preview_slides)} 页 PPT（主题：{theme_colors['name']}{img_info}）\n"
            f"下载链接：{download_url}"
        ),
    }
